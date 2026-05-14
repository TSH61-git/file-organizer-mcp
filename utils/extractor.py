import os
import PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation
from bs4 import BeautifulSoup

class FileContentExtractor:
    @staticmethod
    def get_text(file_path: str, max_length: int) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        
        try:
            # PDF
            if ext == '.pdf':
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text = ""
                    for page in reader.pages[:2]: # קורא עד 2 עמודים ראשונים לתצוגה מקדימה
                        text += page.extract_text() + " "
                    return text[:max_length]

            # Word
            elif ext == '.docx':
                doc = Document(file_path)
                return "\n".join([p.text for p in doc.paragraphs])[:max_length]

            # Excel
            elif ext == '.xlsx':
                wb = openpyxl.load_workbook(file_path, read_only=True)
                sheet = wb.active
                data = []
                for row in sheet.iter_rows(max_row=10, values_only=True):
                    data.append(" | ".join([str(c) for c in row if c is not None]))
                return "\n".join(data)[:max_length]

            # PowerPoint
            elif ext == '.pptx':
                prs = Presentation(file_path)
                slides_content = []
                
                # סורק עד 5 שקופיות ראשונות לקבלת הקשר רחב יותר
                for i, slide in enumerate(prs.slides[:5]):
                    slide_text = []
                    
                    # ניסיון לחלץ קודם כל את כותרת השקף (Title)
                    if slide.shapes.title:
                        slide_text.append(f"[כותרת שקף {i+1}: {slide.shapes.title.text.strip()}]")
                    
                    # חילוץ שאר הטקסט מהשקף
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            # נמנע מהכפלת הכותרת אם כבר הוספנו אותה
                            if slide.shapes.title and shape == slide.shapes.title:
                                continue
                            slide_text.append(shape.text.strip())
                    
                    if slide_text:
                        slides_content.append(" ".join(slide_text))
                
                return "\n".join(slides_content)[:max_length]
            
            # HTML - חילוץ חכם
            elif ext in ['.html', '.htm']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    soup = BeautifulSoup(f, 'html.parser')
                    
                    # הסרת אלמנטים לא טקסטואליים
                    for script_or_style in soup(["script", "style", "meta", "noscript", "header", "footer"]):
                        script_or_style.decompose()
                    
                    # חילוץ טקסט נקי
                    text = soup.get_text(separator=' ')
                    
                    # ניקוי רווחים כפולים וירידות שורה מיותרות
                    lines = (line.strip() for line in text.splitlines())
                    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                    clean_text = '\n'.join(chunk for chunk in chunks if chunk)
                    
                    return clean_text[:max_length]

            # טקסט רגיל, CSV, קוד
            elif ext in ['.txt', '.csv', '.py', '.json', '.html']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read(max_length)

            return f"פורמט {ext} נסרק (לא חולץ טקסט)"

        except Exception as e:
            return f"שגיאה בחילוץ תוכן: {str(e)}"